#_*_ encoding: utf8 _*_
import os
from datetime import datetime

import copy
from pptx import Presentation
import win32com.client

# 저장할 위치 가져오기
path = r"\\***\fileserver\***\주간업무보고"
employee_path = path + "\\" + "주간보고 인원.txt"

# 직원 목록 가져오기
f = open(employee_path, "r", encoding='UTF8')
employee_list = f.readlines()

for i, line in enumerate(employee_list):
    employee_list[i] = line.strip()
f.close()
print("Employee: ", employee_list)

# 오늘 날짜 가져오기
y = datetime.today().year
yy = (datetime.today().year) % 100
m = datetime.today().strftime("%m")
d = datetime.today().strftime("%d")

today = str(yy) + str(m) + str(d)

# 이번주 폴더 열기
path = path + "\\" + str(y) + "-주간보고" + "\\" + today
employee = os.listdir(path) # return type: list

# 이번주 통합 ppt 제목 (yyyymmdd_주간업무보고)
total_ppt_name = str(y) + str(m) + str(d) + "_주간업무보고.pptx"

# presentation 객체 생성
prs = None

# #슬라이드 레이아웃 선택
# slide_layout = prs.slide_layouts[0]

#차례대로 가져오기
for ql in employee_list:
    tmp = [s for s in employee if ql in s]
    if tmp:
        personal_ppt_name = tmp[0]
    else:
        continue
    personal_ppt_path = path + "\\" + personal_ppt_name

    if personal_ppt_path[-3:] == "ppt":
        tmp_name = personal_ppt_path + "x"
        if tmp_name not in employee_list:
            PptApp = win32com.client.Dispatch("Powerpoint.Application")
            PptApp.Visible = True
            PPtPresentation = PptApp.Presentations.Open(personal_ppt_path)

            PPtPresentation.SaveAs(tmp_name, 24)
            PPtPresentation.Close()
            personal_ppt_path = tmp_name
            PptApp.Quit()
        else:
            continue
        
    print(personal_ppt_path)
    if prs == None:
        prs = Presentation(personal_ppt_path)
        #슬라이드 레이아웃 선택
        slide_layout = prs.slide_layouts[0]

    else:
        #선택한 레이아웃의 슬라이드 추가
        copy_slide = prs.slides.add_slide(slide_layout)

        copy_prs = Presentation(personal_ppt_path)
        source_slide = copy_prs.slides[0]
        
        for shape in source_slide.shapes:
            el = shape.element
            newel = copy.deepcopy(el)
            copy_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')


save_ppt_name = path + "\\" + total_ppt_name
prs.save(save_ppt_name)
